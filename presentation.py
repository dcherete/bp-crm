"""
Projeto Sovereign — Apresentacao Financeira Estrategica
Brasil Paralelo | Abril 2026

Gera: projeto-sovereign-apresentacao.pptx
Dependencia: python-pptx

Uso: python presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Pt
import copy

# ---------------------------------------------------------------------------
# PALETA DE CORES — sistema narrativo
# ---------------------------------------------------------------------------
DARK_BG       = RGBColor(0x0D, 0x11, 0x17)   # fundo — autoridade, confianca
CARD_BG       = RGBColor(0x16, 0x1B, 0x22)   # superficie elevada
SURFACE       = RGBColor(0x21, 0x26, 0x2D)   # elementos interativos
BORDER        = RGBColor(0x30, 0x36, 0x3D)   # divisores suaves

RED_CHALLENGE = RGBColor(0xF8, 0x51, 0x49)   # desafio, problema, risco
AMBER_PIVOT   = RGBColor(0xE3, 0xB3, 0x41)   # aviso, ponto de decisao
GREEN_GROWTH  = RGBColor(0x3F, 0xB9, 0x50)   # crescimento, economia, positivo
BLUE_DATA     = RGBColor(0x58, 0xA6, 0xFF)   # dado neutro, informacao
PURPLE_INSIGHT= RGBColor(0xBC, 0x8C, 0xFF)   # insight, perspectiva estrategica
WHITE_PRIMARY = RGBColor(0xF0, 0xF6, 0xFC)   # texto principal
GRAY_SECONDARY= RGBColor(0x8B, 0x94, 0x9E)   # texto secundario

# Hexadecimais para fill direto
HEX_DARK_BG        = "0D1117"
HEX_CARD_BG        = "161B22"
HEX_SURFACE        = "21262D"
HEX_BORDER         = "30363D"
HEX_RED            = "F85149"
HEX_AMBER          = "E3B341"
HEX_GREEN          = "3FB950"
HEX_BLUE           = "58A6FF"
HEX_PURPLE         = "BC8CFF"
HEX_WHITE          = "F0F6FC"
HEX_GRAY           = "8B949E"

# ---------------------------------------------------------------------------
# DIMENSOES — 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# HELPERS GERAIS
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # layout completamente em branco
    slide  = prs.slides.add_slide(layout)
    return slide


def set_bg(slide, color_hex=HEX_DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)


def add_rect(slide, left, top, width, height, fill_hex, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE_PRIMARY,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_paragraph(tf, text, font_size=16, bold=False, color=WHITE_PRIMARY,
                  align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p   = tf.add_paragraph()
    p.alignment    = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None,
                 title_color=WHITE_PRIMARY, sub_color=BLUE_DATA):
    # Faixa superior discreta
    add_rect(slide, 0, 0, 13.33, 0.06, HEX_BORDER)
    add_text_box(slide, title,
                 left=0.5, top=0.18, width=12.3, height=0.8,
                 font_size=40, bold=True, color=title_color,
                 align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     left=0.5, top=0.95, width=12.3, height=0.45,
                     font_size=22, bold=False, color=sub_color,
                     align=PP_ALIGN.LEFT)


def slide_number(slide, num, total=12):
    label = f"{num} / {total}"
    add_text_box(slide, label,
                 left=12.2, top=7.1, width=1.0, height=0.3,
                 font_size=11, color=GRAY_SECONDARY, align=PP_ALIGN.RIGHT)


def brand_footer(slide):
    add_rect(slide, 0, 7.28, 13.33, 0.22, HEX_CARD_BG)
    add_text_box(slide, "Brasil Paralelo  |  Projeto Sovereign  |  Abril 2026  |  CONFIDENCIAL",
                 left=0.4, top=7.3, width=12.5, height=0.18,
                 font_size=10, color=GRAY_SECONDARY, align=PP_ALIGN.LEFT)


def card(slide, left, top, width, height, fill_hex=HEX_CARD_BG, border_hex=HEX_BORDER):
    bg = add_rect(slide, left, top, width, height, fill_hex)
    border_shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    border_shape.fill.background()
    border_shape.line.color.rgb = RGBColor.from_string(border_hex)
    border_shape.line.width = Pt(0.75)
    return bg


def hero_number(slide, value, label, cx, cy,
                num_color=GREEN_GROWTH, label_color=GRAY_SECONDARY,
                num_size=56, lbl_size=13):
    add_text_box(slide, value,
                 left=cx - 2.5, top=cy, width=5, height=1.1,
                 font_size=num_size, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, label,
                 left=cx - 2.5, top=cy + 1.0, width=5, height=0.4,
                 font_size=lbl_size, color=label_color,
                 align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# SLIDE 1 — Hook: "R$ 10,9 milhoes por ano. Existe uma alternativa?"
# ---------------------------------------------------------------------------

def slide_01(prs):
    slide = blank_slide(prs)
    set_bg(slide)

    # Linha de acento lateral esquerda
    add_rect(slide, 0, 0, 0.08, 7.5, HEX_AMBER)

    # Etiqueta acima do titulo
    add_text_box(slide, "APRESENTACAO ESTRATEGICA  |  ABRIL 2026",
                 left=0.6, top=1.5, width=12, height=0.35,
                 font_size=12, color=GRAY_SECONDARY, align=PP_ALIGN.LEFT)

    # Titulo principal — a pergunta que abre a narrativa
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.0), Inches(2.2))
    tf  = txb.text_frame
    tf.word_wrap = True
    p1  = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1  = p1.add_run()
    r1.text = "R$ 10,9 milhoes por ano."
    r1.font.size  = Pt(52)
    r1.font.bold  = True
    r1.font.color.rgb = WHITE_PRIMARY

    add_paragraph(tf, "Existe uma alternativa?",
                  font_size=52, bold=True, color=AMBER_PIVOT,
                  align=PP_ALIGN.LEFT, space_before=2)

    # Subtitulo descritivo
    add_text_box(slide,
                 "O contrato Insider vence em setembro de 2026. Uma contra-proposta chegou. "
                 "E construimos nossa propria stack. Este e o momento de decidir.",
                 left=0.6, top=4.4, width=9.5, height=0.9,
                 font_size=17, color=GRAY_SECONDARY, align=PP_ALIGN.LEFT)

    # Metricas de contexto rapidas
    metrics = [
        ("POC", "Concluida\nmar/2026"),
        ("Throughput", "372 jobs/s\nvalidado"),
        ("Deadline", "14/set/2026"),
    ]
    for i, (lbl, val) in enumerate(metrics):
        x = 0.6 + i * 3.2
        card(slide, x, 5.55, 2.9, 1.5)
        add_text_box(slide, lbl,
                     left=x + 0.15, top=5.65, width=2.6, height=0.35,
                     font_size=11, color=GRAY_SECONDARY)
        add_text_box(slide, val,
                     left=x + 0.15, top=5.98, width=2.6, height=0.7,
                     font_size=15, bold=True, color=WHITE_PRIMARY)

    brand_footer(slide)
    slide_number(slide, 1)


# ---------------------------------------------------------------------------
# SLIDE 2 — Realidade atual: decomposicao do contrato Insider
# ---------------------------------------------------------------------------

def slide_02(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "O contrato atual: R$ 10,9M/ano",
                 "WhatsApp representa 66% do custo total — e e a variavel que tudo decide.",
                 title_color=WHITE_PRIMARY, sub_color=RED_CHALLENGE)

    # Grafico de barras horizontais mostrando dominancia do WhatsApp
    chart_data = ChartData()
    chart_data.categories = [
        "WhatsApp",
        "Email",
        "Jornadas (Architect)",
        "WebSuite (In-app)",
        "SMS",
        "Push Mobile",
        "Upsert API",
    ]
    chart_data.add_series("Custo Anual (R$)", (
        7_130_025,
        1_063_957,
        441_371,
        441_442,
        348_165,
        101_430,
        8_308,
    ))

    from pptx.enum.chart import XL_CHART_TYPE
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(1.55), Inches(7.8), Inches(5.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Estilo do grafico
    plot = chart.plots[0]
    series = plot.series[0]
    # Barras vermelhas para reforcar o "desafio"
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = RED_CHALLENGE

    chart.chart_title.has_text_frame = False

    # Painel lateral de resumo
    card(slide, 8.5, 1.55, 4.5, 5.5)

    items = [
        ("WhatsApp",         "R$ 7.130.025", HEX_RED,    66.0),
        ("Email",            "R$ 1.063.957", HEX_AMBER,   9.8),
        ("Jornadas",         "R$   441.371", HEX_BLUE,    4.1),
        ("WebSuite",         "R$   441.442", HEX_BLUE,    4.1),
        ("SMS",              "R$   348.165", HEX_GRAY,    3.2),
        ("Push Mobile",      "R$   101.430", HEX_GRAY,    0.9),
        ("Upsert API",       "R$     8.308", HEX_GRAY,    0.1),
    ]

    y = 1.75
    for name, value, color_hex, pct in items:
        # mini barra de progresso
        bar_w_max = 3.6
        bar_w = bar_w_max * (pct / 66.0) if pct > 0 else 0.05
        bar_w = min(bar_w, bar_w_max)
        add_rect(slide, 8.65, y + 0.02, bar_w, 0.22, color_hex)
        add_text_box(slide, name,
                     left=8.65, top=y + 0.26, width=2.2, height=0.25,
                     font_size=11, color=WHITE_PRIMARY)
        add_text_box(slide, value,
                     left=11.0, top=y + 0.26, width=1.8, height=0.25,
                     font_size=11, bold=True,
                     color=RGBColor.from_string(color_hex),
                     align=PP_ALIGN.RIGHT)
        y += 0.72

    # Total destacado
    add_rect(slide, 8.5, 6.7, 4.5, 0.05, HEX_BORDER)
    add_text_box(slide, "TOTAL (com impostos)",
                 left=8.65, top=6.78, width=2.5, height=0.3,
                 font_size=12, color=GRAY_SECONDARY)
    add_text_box(slide, "R$ 10.854.708",
                 left=8.65, top=7.0, width=4.1, height=0.3,
                 font_size=16, bold=True, color=RED_CHALLENGE,
                 align=PP_ALIGN.RIGHT)

    brand_footer(slide)
    slide_number(slide, 2)


# ---------------------------------------------------------------------------
# SLIDE 3 — Ponto de inflexao: dois caminhos
# ---------------------------------------------------------------------------

def slide_03(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Setembro de 2026: dois caminhos",
                 "O contrato vence. A contra-proposta chegou. A stack esta pronta. E hora de escolher.",
                 title_color=WHITE_PRIMARY, sub_color=AMBER_PIVOT)

    # Caminho A — Contra-proposta
    card(slide, 0.4, 1.6, 5.9, 5.3, fill_hex=HEX_CARD_BG, border_hex=HEX_AMBER)
    add_text_box(slide, "CAMINHO A",
                 left=0.6, top=1.75, width=5.5, height=0.35,
                 font_size=11, color=AMBER_PIVOT, align=PP_ALIGN.LEFT)
    add_text_box(slide, "Aceitar a\nContra-proposta Insider",
                 left=0.6, top=2.1, width=5.5, height=0.85,
                 font_size=24, bold=True, color=WHITE_PRIMARY)

    pontos_a = [
        ("+ Economia imediata de R$ 2,25M/ano", HEX_GREEN),
        ("+ Push ilimitado e gratuito", HEX_GREEN),
        ("+ Preco em BRL (sem risco cambial)", HEX_GREEN),
        ("+ Suporte e SLA incluidos", HEX_GREEN),
        ("- Lock-in de 24 meses", HEX_RED),
        ("- R$ 540.000 upfront imediatos", HEX_RED),
        ("- Dependencia continua de terceiro", HEX_RED),
        ("- Em set/2028 o ciclo se repete", HEX_RED),
    ]
    y = 3.1
    for txt, color_hex in pontos_a:
        add_text_box(slide, txt,
                     left=0.65, top=y, width=5.4, height=0.28,
                     font_size=13,
                     color=RGBColor.from_string(color_hex))
        y += 0.33

    add_text_box(slide, "R$ 8.605.572/ano",
                 left=0.6, top=6.3, width=5.5, height=0.4,
                 font_size=20, bold=True, color=AMBER_PIVOT,
                 align=PP_ALIGN.LEFT)

    # Divisor central
    add_rect(slide, 6.42, 1.6, 0.06, 5.3, HEX_BORDER)
    add_text_box(slide, "OU",
                 left=6.1, top=4.0, width=0.7, height=0.45,
                 font_size=18, bold=True, color=GRAY_SECONDARY,
                 align=PP_ALIGN.CENTER)

    # Caminho B — Stack interna
    card(slide, 6.6, 1.6, 6.35, 5.3, fill_hex=HEX_CARD_BG, border_hex=HEX_BLUE)
    add_text_box(slide, "CAMINHO B",
                 left=6.8, top=1.75, width=6.0, height=0.35,
                 font_size=11, color=BLUE_DATA, align=PP_ALIGN.LEFT)
    add_text_box(slide, "Implantar a\nStack Interna Sovereign",
                 left=6.8, top=2.1, width=6.0, height=0.85,
                 font_size=24, bold=True, color=WHITE_PRIMARY)

    pontos_b = [
        ("+ Sem lock-in — liberdade total", HEX_GREEN),
        ("+ Controle total dos dados", HEX_GREEN),
        ("+ Sem desembolso upfront", HEX_GREEN),
        ("+ Integracao BigQuery/modelos preditivos", HEX_GREEN),
        ("+ SMS eliminado (R$ 272K/ano a menos)", HEX_GREEN),
        ("- Custo pode ser maior se WA ratio < 1,42", HEX_AMBER),
        ("- Equipe de infra dedicada necessaria", HEX_AMBER),
        ("- Responsabilidade operacional interna", HEX_AMBER),
    ]
    y = 3.1
    for txt, color_hex in pontos_b:
        add_text_box(slide, txt,
                     left=6.85, top=y, width=5.9, height=0.28,
                     font_size=13,
                     color=RGBColor.from_string(color_hex))
        y += 0.33

    add_text_box(slide, "R$ 9.055.440 – 9.480.384/ano",
                 left=6.8, top=6.3, width=5.9, height=0.4,
                 font_size=18, bold=True, color=BLUE_DATA,
                 align=PP_ALIGN.LEFT)

    brand_footer(slide)
    slide_number(slide, 3)


# ---------------------------------------------------------------------------
# SLIDE 4 — Tres cenarios: comparacao visual
# ---------------------------------------------------------------------------

def slide_04(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Tres cenarios: visao geral",
                 "Comparacao de custo anual (com impostos) — valores em R$",
                 title_color=WHITE_PRIMARY, sub_color=BLUE_DATA)

    chart_data = ChartData()
    chart_data.categories = ["Contrato Atual", "Contra-proposta\nInsider", "Stack Interna\n(WA 1:1)", "Stack Interna\n(WA 1,5 avg)"]
    chart_data.add_series("Custo Anual (R$)", (
        10_854_708,
        8_605_572,
        9_267_912,   # media do range 9.055.440–9.480.384
        6_218_676,
    ))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.55), Inches(7.8), Inches(5.4),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = False

    # Cores individuais por barra (workaround via series fill)
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BLUE_DATA


    # Cards de insight lateral
    insights = [
        (HEX_RED,    "Contrato Atual",        "R$ 10.854.708",  "Base de referencia — custo pleno"),
        (HEX_AMBER,  "Contra-proposta",       "R$ 8.605.572",   "-21% vs atual | lock-in 24m"),
        (HEX_BLUE,   "Stack Interna (1:1 WA)","R$ 9.267.912",   "-15% vs atual | sem lock-in"),
        (HEX_GREEN,  "Stack Interna (1,5 WA)","R$ 6.218.676",   "-43% vs atual | cenario otimo"),
    ]

    y = 1.55
    for color_hex, name, value, note in insights:
        card(slide, 8.5, y, 4.5, 1.2, fill_hex=HEX_CARD_BG)
        add_rect(slide, 8.5, y, 0.07, 1.2, color_hex)
        add_text_box(slide, name,
                     left=8.65, top=y + 0.1, width=4.2, height=0.3,
                     font_size=12, color=GRAY_SECONDARY)
        add_text_box(slide, value,
                     left=8.65, top=y + 0.38, width=4.2, height=0.42,
                     font_size=22, bold=True,
                     color=RGBColor.from_string(color_hex))
        add_text_box(slide, note,
                     left=8.65, top=y + 0.85, width=4.2, height=0.25,
                     font_size=11, color=GRAY_SECONDARY)
        y += 1.35

    brand_footer(slide)
    slide_number(slide, 4)


# ---------------------------------------------------------------------------
# SLIDE 5 — Contra-proposta Insider: o que voce ganha, o que custa
# ---------------------------------------------------------------------------

def slide_05(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Contra-proposta Insider: o que esta em jogo",
                 "24 meses de lock-in. R$ 540K na mesa. E uma reducao real de 21%.",
                 title_color=WHITE_PRIMARY, sub_color=AMBER_PIVOT)

    # Bloco de economia
    card(slide, 0.4, 1.55, 4.2, 5.5)
    add_text_box(slide, "ECONOMIA ANUAL",
                 left=0.6, top=1.7, width=3.8, height=0.35,
                 font_size=11, color=GRAY_SECONDARY)
    add_text_box(slide, "-R$ 2.249.136",
                 left=0.6, top=2.05, width=3.8, height=0.75,
                 font_size=36, bold=True, color=GREEN_GROWTH)
    add_text_box(slide, "-21% vs contrato atual",
                 left=0.6, top=2.75, width=3.8, height=0.3,
                 font_size=14, color=GRAY_SECONDARY)

    add_rect(slide, 0.6, 3.2, 3.8, 0.04, HEX_BORDER)

    details = [
        ("Custo anual com impostos",     "R$ 8.605.572", HEX_AMBER),
        ("Custo mensal medio",           "R$ 717.131",   HEX_WHITE),
        ("Pagamento upfront",            "R$ 540.000",   HEX_RED),
        ("WhatsApp por disparo",         "R$ 0,3726",    HEX_WHITE),
        ("Push (mobile + web)",          "GRATUITO",     HEX_GREEN),
        ("Duracao do contrato",          "24 meses",     HEX_AMBER),
        ("Risco cambial",                "Nenhum (BRL)",  HEX_GREEN),
    ]
    y = 3.35
    for lbl, val, color_hex in details:
        add_text_box(slide, lbl,
                     left=0.6, top=y, width=2.5, height=0.28,
                     font_size=12, color=GRAY_SECONDARY)
        add_text_box(slide, val,
                     left=3.0, top=y, width=1.5, height=0.28,
                     font_size=12, bold=True,
                     color=RGBColor.from_string(color_hex),
                     align=PP_ALIGN.RIGHT)
        y += 0.38

    # Bloco de lock-in — o "custo oculto"
    card(slide, 4.85, 1.55, 8.1, 2.5, border_hex=HEX_RED)
    add_rect(slide, 4.85, 1.55, 8.1, 0.06, HEX_RED)
    add_text_box(slide, "O CUSTO REAL: O LOCK-IN",
                 left=5.05, top=1.65, width=7.7, height=0.35,
                 font_size=12, color=RED_CHALLENGE, bold=True)
    add_text_box(slide,
                 "Ao assinar, a Brasil Paralelo perde a capacidade de negociar "
                 "por 24 meses. Em setembro de 2028, o ciclo se repete — "
                 "com a Insider ainda mais entrincheirada. O R$ 540K upfront "
                 "nao e desconto: e um sinal de que voce nao vai sair.",
                 left=5.05, top=2.05, width=7.7, height=1.7,
                 font_size=14, color=WHITE_PRIMARY)

    # Bloco de linha do tempo
    card(slide, 4.85, 4.3, 8.1, 2.75)
    add_text_box(slide, "LINHA DO TEMPO DO CONTRATO",
                 left=5.05, top=4.42, width=7.7, height=0.3,
                 font_size=11, color=GRAY_SECONDARY)

    timeline = [
        (0.3,  "Abr/2026",    "Assinatura\n+ R$540K",  HEX_RED),
        (2.05, "Set/2026",    "Stack Sovereign\nconcluida",  HEX_BLUE),
        (3.8,  "Mar/2028",    "Meio do\nlock-in",       HEX_AMBER),
        (5.55, "Set/2028",    "NOVO CICLO\nde negociacao", HEX_AMBER),
    ]
    for x_offset, date, label, color_hex in timeline:
        x = 5.1 + x_offset
        add_rect(slide, x, 4.85, 0.12, 0.12, color_hex)
        add_text_box(slide, date,
                     left=x - 0.3, top=5.05, width=1.0, height=0.25,
                     font_size=10, color=GRAY_SECONDARY, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     left=x - 0.4, top=5.3, width=1.2, height=0.5,
                     font_size=10, bold=True,
                     color=RGBColor.from_string(color_hex),
                     align=PP_ALIGN.CENTER)

    # Linha da timeline
    add_rect(slide, 5.1, 4.9, 7.6, 0.03, HEX_BORDER)

    brand_footer(slide)
    slide_number(slide, 5)


# ---------------------------------------------------------------------------
# SLIDE 6 — Stack interna: o que construimos
# ---------------------------------------------------------------------------

def slide_06(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Stack Sovereign: o que construimos",
                 "POC concluida. Throughput validado. Custo detalhado por componente.",
                 title_color=WHITE_PRIMARY, sub_color=BLUE_DATA)

    # Grafico de barras dos componentes de custo
    chart_data = ChartData()
    chart_data.categories = [
        "WhatsApp\n(Meta API)",
        "Amazon\nSES",
        "PostHog\nCloud",
        "MongoDB\nAtlas",
        "Infra AWS\n(ECS/Redis)",
        "Firebase\nFCM",
    ]
    chart_data.add_series("Custo Mensal (R$)", (
        660_000,
        66_987,
        37_982,   # media do range
        2_962,    # media do range
        4_894,
        0,
    ))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.55), Inches(7.0), Inches(5.4),
        chart_data
    )
    ch = chart_frame.chart
    ch.has_legend = False
    ch.chart_title.has_text_frame = False
    ch.plots[0].series[0].format.fill.solid()
    ch.plots[0].series[0].format.fill.fore_color.rgb = BLUE_DATA

    # Tabela de custos detalhados
    card(slide, 7.6, 1.55, 5.35, 5.4)
    add_text_box(slide, "CUSTO MENSAL DETALHADO",
                 left=7.8, top=1.68, width=5.0, height=0.3,
                 font_size=11, color=GRAY_SECONDARY)

    rows = [
        ("Amazon SES (104M emails)",    "R$ 66.987",   HEX_AMBER),
        ("Infra AWS (ECS, Redis, LB)",  "R$ 4.894",    HEX_AMBER),
        ("MongoDB Atlas M30",           "R$ 2.857–3.068", HEX_BLUE),
        ("PostHog Cloud (37M eventos)", "R$ 19.882–56.083", HEX_BLUE),
        ("Firebase FCM (push)",         "GRATUITO",    HEX_GREEN),
        ("WhatsApp Meta API (1,25M)",   "R$ 660.000",  HEX_RED),
        ("SMS",                         "R$ 0 (eliminado)", HEX_GREEN),
    ]
    y = 2.08
    for lbl, val, color_hex in rows:
        add_text_box(slide, lbl,
                     left=7.8, top=y, width=3.3, height=0.28,
                     font_size=12, color=WHITE_PRIMARY)
        add_text_box(slide, val,
                     left=11.0, top=y, width=1.7, height=0.28,
                     font_size=12, bold=True,
                     color=RGBColor.from_string(color_hex),
                     align=PP_ALIGN.RIGHT)
        y += 0.44

    add_rect(slide, 7.8, 5.2, 4.9, 0.04, HEX_BORDER)
    add_text_box(slide, "Total anual (range)",
                 left=7.8, top=5.28, width=3.3, height=0.3,
                 font_size=13, color=GRAY_SECONDARY)
    add_text_box(slide, "R$ 9,05M – 9,48M",
                 left=7.8, top=5.6, width=4.9, height=0.42,
                 font_size=20, bold=True, color=BLUE_DATA,
                 align=PP_ALIGN.RIGHT)

    # Badge de POC validada
    add_rect(slide, 7.8, 6.25, 4.9, 0.6, HEX_SURFACE)
    add_text_box(slide,
                 "POC validada: 372 jobs/s  |  Necessario: 280/s no pico  |  Margem: +33%",
                 left=7.85, top=6.3, width=4.8, height=0.45,
                 font_size=11, color=GREEN_GROWTH, align=PP_ALIGN.CENTER)

    brand_footer(slide)
    slide_number(slide, 6)


# ---------------------------------------------------------------------------
# SLIDE 7 — A variavel WhatsApp: o climax da decisao
# ---------------------------------------------------------------------------

def slide_07(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "A variavel que decide tudo: WhatsApp",
                 "Insider cobra por DISPARO. Meta API cobra por CONVERSA de 24h. O breakeven esta em 1,42 msgs/conversa.",
                 title_color=WHITE_PRIMARY, sub_color=PURPLE_INSIGHT)

    # Grafico de linha — curva de custo por razao msgs/conversa
    chart_data = ChartData()
    chart_data.categories = ["1,0", "1,2", "1,42", "1,6", "2,0", "2,5", "3,0"]

    # Custo anual da stack interna varia com o ratio
    # Base: 1,25M conversas/mes * 12 = 15M/ano
    # Insider: 15M * R$0,3726 = R$5.589.000/ano (so WA)
    # Meta API: custo por conversa = R$0,528 (com fator 1,4832)
    # Stack interna total = custos fixos (~R$1.074.000/ano) + Meta API WA
    custos_fixos_anuais = 1_074_000  # SES + infra + PostHog + MongoDB
    meta_por_conversa   = 0.528
    conversas_anuais    = 15_000_000
    insider_wa_por_envio = 0.3726

    ratios = [1.0, 1.2, 1.42, 1.6, 2.0, 2.5, 3.0]

    # Stack interna: paga por conversa, nao por mensagem
    # Com ratio = N: sao 15M/N conversas para 15M msgs equivalentes
    # Mas na pratica, com ratio > 1, voce envia MAIS msgs pelo mesmo custo
    # Para comparar justo: numero de conversas = dispatches / ratio
    # Custo WA stack = (15M / ratio) * 0.528

    stack_costs  = [custos_fixos_anuais + (conversas_anuais / r) * meta_por_conversa for r in ratios]
    insider_cost = 8_605_572  # contra-proposta anual total

    insider_line = [insider_cost] * len(ratios)

    chart_data.add_series("Stack Interna (custo anual)", [round(c) for c in stack_costs])
    chart_data.add_series("Contra-proposta Insider",     insider_line)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.4), Inches(1.55), Inches(8.3), Inches(5.4),
        chart_data
    )
    ch = chart_frame.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.chart_title.has_text_frame = False

    # Linha da stack interna — azul
    s0 = ch.plots[0].series[0]
    s0.format.line.color.rgb = BLUE_DATA
    s0.format.line.width = Pt(2.5)

    # Linha da Insider — amarelo
    s1 = ch.plots[0].series[1]
    s1.format.line.color.rgb = AMBER_PIVOT
    s1.format.line.width = Pt(2.5)


    # Painel lateral — analise do breakeven
    card(slide, 9.0, 1.55, 4.0, 5.4, border_hex=HEX_PURPLE)
    add_rect(slide, 9.0, 1.55, 4.0, 0.06, HEX_PURPLE)

    add_text_box(slide, "PONTO DE BREAKEVEN",
                 left=9.15, top=1.65, width=3.7, height=0.3,
                 font_size=11, color=PURPLE_INSIGHT, bold=True)

    add_text_box(slide, "1,42",
                 left=9.15, top=2.0, width=3.7, height=0.9,
                 font_size=64, bold=True, color=PURPLE_INSIGHT,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "mensagens por conversa",
                 left=9.15, top=2.85, width=3.7, height=0.3,
                 font_size=13, color=GRAY_SECONDARY,
                 align=PP_ALIGN.CENTER)

    add_rect(slide, 9.15, 3.3, 3.7, 0.04, HEX_BORDER)

    add_text_box(slide, "Custo Insider (por disparo):",
                 left=9.15, top=3.42, width=3.7, height=0.25,
                 font_size=12, color=GRAY_SECONDARY)
    add_text_box(slide, "R$ 0,3726/msg",
                 left=9.15, top=3.65, width=3.7, height=0.3,
                 font_size=14, bold=True, color=AMBER_PIVOT)

    add_text_box(slide, "Custo Meta API (por conversa):",
                 left=9.15, top=4.05, width=3.7, height=0.25,
                 font_size=12, color=GRAY_SECONDARY)
    add_text_box(slide, "R$ 0,528/conv",
                 left=9.15, top=4.28, width=3.7, height=0.3,
                 font_size=14, bold=True, color=BLUE_DATA)

    add_rect(slide, 9.15, 4.75, 3.7, 0.04, HEX_BORDER)

    add_text_box(slide, "Broadcasts puros:",
                 left=9.15, top=4.85, width=3.7, height=0.25,
                 font_size=11, color=GRAY_SECONDARY)
    add_text_box(slide, "1,0 msg/conv = Insider mais barato",
                 left=9.15, top=5.08, width=3.7, height=0.3,
                 font_size=11, color=RED_CHALLENGE)

    add_text_box(slide, "Sequencias no dia / dialogos:",
                 left=9.15, top=5.5, width=3.7, height=0.25,
                 font_size=11, color=GRAY_SECONDARY)
    add_text_box(slide, "2,0+ msg/conv = Stack muito mais barata",
                 left=9.15, top=5.73, width=3.7, height=0.3,
                 font_size=11, color=GREEN_GROWTH)

    add_text_box(slide, "Taxa de resposta WA Business: 35%\nTemplates interativos: 45–60%",
                 left=9.15, top=6.25, width=3.7, height=0.5,
                 font_size=10, color=GRAY_SECONDARY, italic=True)

    brand_footer(slide)
    slide_number(slide, 7)


# ---------------------------------------------------------------------------
# SLIDE 8 — Projecao de 24 meses
# ---------------------------------------------------------------------------

def slide_08(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Projecao de 24 meses: onde cada caminho leva",
                 "Total acumulado em 24 meses (com pagamento upfront incluido onde aplicavel)",
                 title_color=WHITE_PRIMARY, sub_color=BLUE_DATA)

    chart_data = ChartData()
    chart_data.categories = [f"M{i+1}" for i in range(24)]

    # Contrato atual: R$904.559/mes acumulado
    atual_cum   = [904_559 * (i+1) for i in range(24)]
    # Contra-proposta: R$540K upfront + R$717.131/mes
    cp_cum      = [540_000 + 717_131 * (i+1) for i in range(24)]
    # Stack interna (media): R$772.000/mes sem upfront
    stack_cum   = [771_287 * (i+1) for i in range(24)]

    chart_data.add_series("Contrato Atual",          atual_cum)
    chart_data.add_series("Contra-proposta Insider", cp_cum)
    chart_data.add_series("Stack Interna (media)",   stack_cum)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.4), Inches(1.55), Inches(9.0), Inches(5.4),
        chart_data
    )
    ch = chart_frame.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.chart_title.has_text_frame = False

    s0 = ch.plots[0].series[0]
    s0.format.line.color.rgb = RED_CHALLENGE
    s0.format.line.width = Pt(2.0)

    s1 = ch.plots[0].series[1]
    s1.format.line.color.rgb = AMBER_PIVOT
    s1.format.line.width = Pt(2.0)

    s2 = ch.plots[0].series[2]
    s2.format.line.color.rgb = BLUE_DATA
    s2.format.line.width = Pt(2.0)


    # Cards de total aos 24 meses
    cards_24 = [
        ("Contrato Atual",          "R$ 21.709.416", HEX_RED,   "Base"),
        ("Contra-proposta",         "R$ 17.751.144", HEX_AMBER, "Inclui R$540K upfront"),
        ("Stack Interna",           "R$ 18.110.880", HEX_BLUE,  "Cenario WA 1:1"),
    ]

    y = 1.55
    for name, total, color_hex, note in cards_24:
        card(slide, 9.6, y, 3.4, 1.6)
        add_rect(slide, 9.6, y, 0.07, 1.6, color_hex)
        add_text_box(slide, name,
                     left=9.75, top=y + 0.1, width=3.1, height=0.3,
                     font_size=11, color=GRAY_SECONDARY)
        add_text_box(slide, total,
                     left=9.75, top=y + 0.4, width=3.1, height=0.5,
                     font_size=19, bold=True,
                     color=RGBColor.from_string(color_hex))
        add_text_box(slide, note,
                     left=9.75, top=y + 0.9, width=3.1, height=0.3,
                     font_size=10, color=GRAY_SECONDARY, italic=True)
        y += 1.75

    # Nota de rodape do slide
    add_text_box(slide,
                 "Nota: Stack Interna pode ser igual ou mais barata que a contra-proposta "
                 "se a media de msgs/conversa no WhatsApp superar 1,42.",
                 left=0.4, top=6.85, width=13.0, height=0.35,
                 font_size=10, color=GRAY_SECONDARY, italic=True)

    brand_footer(slide)
    slide_number(slide, 8)


# ---------------------------------------------------------------------------
# SLIDE 9 — Analise de sensibilidade: o mapa da decisao
# ---------------------------------------------------------------------------

def slide_09(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Analise de sensibilidade: o mapa da decisao",
                 "Diferenca de custo em 24 meses (Stack Interna vs Contra-proposta) por ratio msgs/conversa",
                 title_color=WHITE_PRIMARY, sub_color=PURPLE_INSIGHT)

    # Diferenca: Stack - Contra-proposta em 24 meses
    # Positivo = Stack mais cara | Negativo = Stack mais barata
    chart_data = ChartData()
    chart_data.categories = [
        "1,0 msg/conv\n(broadcast puro)",
        "1,42 msg/conv\n(breakeven)",
        "2,0 msg/conv\n(sequencias)",
        "3,0 msg/conv\n(alto engajamento)",
    ]
    chart_data.add_series("Diferenca em 24 meses (R$)", (
         359_736,   # pior caso — stack mais cara
               0,   # breakeven
        -7_560_264, # stack mais barata
        -10_200_264,# stack muito mais barata
    ))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(1.55), Inches(7.8), Inches(5.4),
        chart_data
    )
    ch = chart_frame.chart
    ch.has_legend = False
    ch.chart_title.has_text_frame = False
    ch.plots[0].series[0].format.fill.solid()
    ch.plots[0].series[0].format.fill.fore_color.rgb = PURPLE_INSIGHT

    # Painel de interpretacao
    card(slide, 8.5, 1.55, 4.5, 5.4)
    add_text_box(slide, "INTERPRETACAO",
                 left=8.68, top=1.68, width=4.1, height=0.3,
                 font_size=11, color=GRAY_SECONDARY, bold=True)

    interps = [
        ("1,0 msg/conv",   "Broadcasts sem resposta. Stack custa R$359K a mais em 24m.",
         HEX_RED, HEX_RED),
        ("1,42 msg/conv",  "Ponto exato de equilibrio. Custo identico ao longo de 24 meses.",
         HEX_AMBER, HEX_AMBER),
        ("2,0 msg/conv",   "Sequencias simples no mesmo dia. Stack economiza R$7,6M em 24m.",
         HEX_GREEN, HEX_GREEN),
        ("3,0 msg/conv",   "Alto engajamento / dialogos. Stack economiza R$10,2M em 24m.",
         HEX_GREEN, HEX_GREEN),
    ]
    y = 2.1
    for ratio, desc, accent_hex, _ in interps:
        add_rect(slide, 8.5, y, 0.06, 1.15, accent_hex)
        add_text_box(slide, ratio,
                     left=8.65, top=y + 0.04, width=3.7, height=0.3,
                     font_size=13, bold=True,
                     color=RGBColor.from_string(accent_hex))
        add_text_box(slide, desc,
                     left=8.65, top=y + 0.34, width=3.7, height=0.65,
                     font_size=11, color=WHITE_PRIMARY)
        y += 1.28

    # Caixa de conclusao
    add_rect(slide, 8.5, 6.65, 4.5, 0.6, HEX_SURFACE)
    add_text_box(slide,
                 "A decisao e empirica: levante o historico de msgs/conversa do WhatsApp BP.",
                 left=8.6, top=6.7, width=4.3, height=0.5,
                 font_size=11, color=PURPLE_INSIGHT, bold=True)

    brand_footer(slide)
    slide_number(slide, 9)


# ---------------------------------------------------------------------------
# SLIDE 10 — Vantagens estrategicas alem do custo
# ---------------------------------------------------------------------------

def slide_10(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Alem do custo: vantagens estrategicas da stack propria",
                 "Independencia, dados, poder de negociacao — ativos que o dinheiro nao captura diretamente.",
                 title_color=WHITE_PRIMARY, sub_color=GREEN_GROWTH)

    advantages = [
        ("Sem lock-in",
         "Liberdade total para ajustar, migrar ou encerrar contratos a qualquer momento. "
         "Nunca mais negocie de uma posicao de dependencia.",
         HEX_GREEN, "SOBERANIA"),

        ("Controle total dos dados",
         "Todos os dados de comportamento e engajamento ficam nos servidores da BP. "
         "Sem transferencia a terceiros. Modelo preditivo proprio viabilizado.",
         HEX_BLUE, "DADOS"),

        ("Eliminacao do SMS",
         "O contrato Insider inclui SMS (R$348K/ano). A stack interna descontinua o canal, "
         "economizando R$272K/ano vs contra-proposta (que mantem SMS).",
         HEX_AMBER, "EFICIENCIA"),

        ("Poder de negociacao futuro",
         "Com a stack pronta como alternativa real, qualquer renegociacao futura com Insider "
         "ou outro fornecedor parte de uma posicao de forca genuina.",
         HEX_PURPLE, "ESTRATEGIA"),

        ("Integracao BigQuery / IA",
         "A arquitetura Novu + PostHog + dados proprios permite integracao direta com "
         "BigQuery e modelos preditivos de churn, LTV e personalizacao avancada.",
         HEX_BLUE, "INOVACAO"),

        ("Independencia tecnologica permanente",
         "A stack e da BP. Cada feature construida e um ativo permanente — nao uma "
         "locacao que precisa ser renovada a cada 24 meses.",
         HEX_GREEN, "LEGADO"),
    ]

    positions = [
        (0.35, 1.6), (4.6, 1.6), (8.85, 1.6),
        (0.35, 4.3), (4.6, 4.3), (8.85, 4.3),
    ]

    for (x, y), (title, desc, color_hex, tag) in zip(positions, advantages):
        card(slide, x, y, 4.0, 2.4)
        add_rect(slide, x, y, 4.0, 0.06, color_hex)
        add_text_box(slide, tag,
                     left=x + 0.15, top=y + 0.12, width=3.7, height=0.25,
                     font_size=10,
                     color=RGBColor.from_string(color_hex))
        add_text_box(slide, title,
                     left=x + 0.15, top=y + 0.38, width=3.7, height=0.4,
                     font_size=15, bold=True, color=WHITE_PRIMARY)
        add_text_box(slide, desc,
                     left=x + 0.15, top=y + 0.82, width=3.7, height=1.4,
                     font_size=11, color=GRAY_SECONDARY)

    brand_footer(slide)
    slide_number(slide, 10)


# ---------------------------------------------------------------------------
# SLIDE 11 — Recomendacao estrategica
# ---------------------------------------------------------------------------

def slide_11(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    slide_header(slide,
                 "Recomendacao estrategica",
                 "A conclusao que os numeros sustentam — independente da decisao final.",
                 title_color=WHITE_PRIMARY, sub_color=GREEN_GROWTH)

    # Caixa principal de recomendacao
    card(slide, 0.4, 1.55, 12.5, 2.2, border_hex=HEX_GREEN)
    add_rect(slide, 0.4, 1.55, 12.5, 0.07, HEX_GREEN)
    add_text_box(slide, "RECOMENDACAO CENTRAL",
                 left=0.6, top=1.67, width=12.0, height=0.3,
                 font_size=11, color=GREEN_GROWTH, bold=True)
    add_text_box(slide,
                 "Continue o deploy da Stack Sovereign, independente de aceitar ou nao a contra-proposta.",
                 left=0.6, top=2.0, width=12.0, height=0.55,
                 font_size=22, bold=True, color=WHITE_PRIMARY)
    add_text_box(slide,
                 "A contra-proposta exige 24 meses — e em setembro/2028 o ciclo se repete. "
                 "Ter a stack pronta como alternativa real e o unico ativo que garante poder de "
                 "negociacao permanente com qualquer fornecedor de CRM.",
                 left=0.6, top=2.6, width=12.0, height=0.85,
                 font_size=14, color=GRAY_SECONDARY)

    # Tres caminhos possiveis
    paths = [
        ("OPCAO 1",
         "Aceitar contra-proposta\n+ continuar deploy",
         "Economize R$2,25M/ano agora. "
         "Mantenha o deploy ativo como plano B. "
         "Em set/2028, escolha sair ou renegociar de forca.",
         HEX_AMBER),
        ("OPCAO 2",
         "Recusar contra-proposta\n+ implantar stack",
         "Assuma a variabilidade do custo WhatsApp. "
         "Ganhe independencia total imediata. "
         "Objetivo: superar 1,42 msgs/conversa.",
         HEX_BLUE),
        ("OPCAO 3",
         "Implantar stack + negociar\nInsider para canais especificos",
         "Stack cuida de email e push. "
         "Insider ou Meta gerencia WhatsApp separado. "
         "Hibrido que maximiza cada canal.",
         HEX_PURPLE),
    ]

    x = 0.4
    for opt, title, desc, color_hex in paths:
        card(slide, x, 4.0, 4.0, 2.85)
        add_rect(slide, x, 4.0, 4.0, 0.07, color_hex)
        add_text_box(slide, opt,
                     left=x + 0.15, top=4.12, width=3.7, height=0.28,
                     font_size=10,
                     color=RGBColor.from_string(color_hex), bold=True)
        add_text_box(slide, title,
                     left=x + 0.15, top=4.42, width=3.7, height=0.65,
                     font_size=15, bold=True, color=WHITE_PRIMARY)
        add_text_box(slide, desc,
                     left=x + 0.15, top=5.1, width=3.7, height=1.55,
                     font_size=12, color=GRAY_SECONDARY)
        x += 4.25

    brand_footer(slide)
    slide_number(slide, 11)


# ---------------------------------------------------------------------------
# SLIDE 12 — Call to action: a unica pergunta que resolve tudo
# ---------------------------------------------------------------------------

def slide_12(prs):
    slide = blank_slide(prs)
    set_bg(slide)

    # Acento lateral
    add_rect(slide, 0, 0, 0.08, 7.5, HEX_PURPLE)

    add_text_box(slide, "PROXIMA ACAO",
                 left=0.6, top=1.3, width=12, height=0.35,
                 font_size=12, color=PURPLE_INSIGHT, bold=True)

    add_text_box(slide,
                 "Qual e a media real de\nmensagens por conversa\nno WhatsApp da BP?",
                 left=0.6, top=1.75, width=10, height=2.4,
                 font_size=46, bold=True, color=WHITE_PRIMARY)

    add_text_box(slide,
                 "Esta unica variavel determina qual caminho economiza mais dinheiro ao longo de 24 meses. "
                 "Extraia do historico operacional: numero de mensagens enviadas / numero de conversas abertas. "
                 "Se o resultado for maior que 1,42 — a Stack Sovereign e claramente a escolha economica.",
                 left=0.6, top=4.35, width=10.5, height=1.6,
                 font_size=16, color=GRAY_SECONDARY)

    # Tres acoes imediatas
    acoes = [
        ("1", "Levantar historico WhatsApp", "Extrair ratio msgs/conversa dos ultimos 90 dias.", HEX_PURPLE),
        ("2", "Calcular custo real",         "Aplicar formula: (15M / ratio) x R$0,528 x 12.", HEX_BLUE),
        ("3", "Decidir com dados",           "Comparar com R$17,75M da contra-proposta e agir.", HEX_GREEN),
    ]

    x = 0.5
    for num, title, desc, color_hex in acoes:
        card(slide, x, 6.05, 4.05, 1.1)
        add_rect(slide, x, 6.05, 0.06, 1.1, color_hex)
        add_text_box(slide, num,
                     left=x + 0.15, top=6.1, width=0.4, height=0.35,
                     font_size=20, bold=True,
                     color=RGBColor.from_string(color_hex),
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, title,
                     left=x + 0.6, top=6.1, width=3.3, height=0.35,
                     font_size=13, bold=True, color=WHITE_PRIMARY)
        add_text_box(slide, desc,
                     left=x + 0.6, top=6.48, width=3.3, height=0.5,
                     font_size=11, color=GRAY_SECONDARY)
        x += 4.28

    # Detalhe de breakeven final
    add_text_box(slide, "Breakeven: 1,42 msgs/conversa",
                 left=11.0, top=2.5, width=2.1, height=1.0,
                 font_size=16, bold=True, color=PURPLE_INSIGHT,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "Deadline: 14/set/2026",
                 left=11.0, top=3.7, width=2.1, height=0.5,
                 font_size=13, color=AMBER_PIVOT,
                 align=PP_ALIGN.CENTER)

    brand_footer(slide)
    slide_number(slide, 12)


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    print("Gerando slides...")
    slide_01(prs); print("  [1/12] Hook — R$ 10,9M/ano")
    slide_02(prs); print("  [2/12] Decomposicao do contrato Insider")
    slide_03(prs); print("  [3/12] Ponto de inflexao — dois caminhos")
    slide_04(prs); print("  [4/12] Tres cenarios — visao geral")
    slide_05(prs); print("  [5/12] Contra-proposta Insider")
    slide_06(prs); print("  [6/12] Stack Sovereign — o que construimos")
    slide_07(prs); print("  [7/12] Variavel WhatsApp — breakeven (climax)")
    slide_08(prs); print("  [8/12] Projecao de 24 meses")
    slide_09(prs); print("  [9/12] Analise de sensibilidade")
    slide_10(prs); print("  [10/12] Vantagens estrategicas")
    slide_11(prs); print("  [11/12] Recomendacao estrategica")
    slide_12(prs); print("  [12/12] Call to action")

    output = "projeto-sovereign-apresentacao.pptx"
    prs.save(output)
    print(f"\nApresentacao gerada: {output}")
    print("Slides: 12  |  Formato: 16:9  |  Tema: dark")


if __name__ == "__main__":
    main()
